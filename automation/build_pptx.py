# automation/build_pptx.py
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

INK=RGBColor(0x49,0x22,0x6F); INK2=RGBColor(0x34,0x15,0x51)
BRAND=RGBColor(0x60,0x2E,0x91)
TEXT=RGBColor(0x35,0x35,0x35); MUTED=RGBColor(0x6E,0x68,0x78)
HAIR=RGBColor(0xE7,0xE1,0xEF); GREEN=RGBColor(0x1E,0x8E,0x3E); AMBER=RGBColor(0xE8,0xA1,0x00)
RED=RGBColor(0xC5,0x22,0x1F); WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xD8,0xD2,0xE2)
F_DISP="Segoe UI Semibold"; F_BODY="Segoe UI"; F_MONO="Consolas"
MESES=["","Janeiro","Fevereiro","Março","Abril","Maio","Junho","Julho","Agosto","Setembro","Outubro","Novembro","Dezembro"]
_SEC = {"n": 0, "total": 0}

def load_facts(path):
    with open(path, encoding="utf-8-sig") as fh: return json.load(fh)

def slug_from_path(path):
    b=os.path.basename(path)
    if b.startswith("facts_"): b=b[6:]
    if b.endswith(".json"): b=b[:-5]
    return b

def month_pt(ym):
    try:
        y,m=ym.split("-"); return "%s de %s"%(MESES[int(m)], y)
    except Exception: return ym

def _g(d,*keys,default=None):
    for k in keys:
        if not isinstance(d,dict): return default
        d=d.get(k)
    return d if d is not None else default

def _pt(n):
    # PT-BR number: thousands '.', decimal ','
    if n is None: return "—"
    if isinstance(n,float) and not n.is_integer():
        s=("%.1f"%n)
    else:
        s=("%d"%int(round(n)))
    a=s.split("."); a[0]="{:,}".format(int(a[0])).replace(",",".")
    return a[0]+(","+a[1] if len(a)>1 else "")

def add_text(container, text, size, font=F_BODY, color=TEXT, bold=False, align=None, spacing=None):
    tf = container if hasattr(container,"paragraphs") else container.text_frame
    p = tf.paragraphs[0] if (len(tf.paragraphs)==1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    if align is not None: p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = font; r.font.bold = bold; r.font.color.rgb = color
    if spacing is not None: r.font._rPr.set("spc", str(int(spacing*100)))
    return p

def _rect(slide,x,y,w,h,fill=None,line=None,line_w=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    s.shadow.inherit=False
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(line_w)
    return s

def _blank(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    return s

def section_header(slide,eyebrow,title):
    _SEC["n"]+=1
    _rect(slide,0,0,0.083,7.5,fill=BRAND)  # spine
    tb=slide.shapes.add_textbox(Inches(0.7),Inches(0.55),Inches(9),Inches(1.1)).text_frame
    tb.word_wrap=True
    add_text(tb,eyebrow.upper(),11,F_MONO,BRAND)
    add_text(tb,title,30,F_DISP,TEXT,bold=True)
    ix=slide.shapes.add_textbox(Inches(11.4),Inches(0.6),Inches(1.6),Inches(0.4)).text_frame
    ix.paragraphs[0].alignment=PP_ALIGN.RIGHT
    add_text(ix,"%02d / %02d"%(_SEC["n"],_SEC["total"]),12,F_MONO,BRAND,align=PP_ALIGN.RIGHT)

def footer(slide,client,mpt):
    _rect(slide,0.7,7.02,11.93,0.012,fill=HAIR)
    l=slide.shapes.add_textbox(Inches(0.7),Inches(7.08),Inches(7),Inches(0.3)).text_frame
    add_text(l,"Contego Security · MSSP Acronis",10,F_MONO,MUTED)
    r=slide.shapes.add_textbox(Inches(8),Inches(7.08),Inches(4.63),Inches(0.3)).text_frame
    r.paragraphs[0].alignment=PP_ALIGN.RIGHT
    add_text(r,"%s · %s · %02d"%(client,mpt,_SEC["n"]),10,F_MONO,MUTED,align=PP_ALIGN.RIGHT)

def posture_meter(slide,score,x,y,w,dark=False):
    # segmented scale: red 0-59, amber 60-84, green 85-100, with a pin at score
    h=0.22
    segs=[(0.59,RED),(0.25,AMBER),(0.16,GREEN)]
    cx=x
    for frac,col in segs:
        _rect(slide,cx,y,w*frac,h,fill=col)
        cx+=w*frac
    px=x+w*(max(0,min(100,score))/100.0)
    _rect(slide,px-0.02,y-0.06,0.04,h+0.12,fill=WHITE if dark else TEXT)
    lab=slide.shapes.add_textbox(Inches(x),Inches(y+h+0.06),Inches(w),Inches(0.3)).text_frame
    lab.paragraphs[0].alignment=PP_ALIGN.LEFT
    add_text(lab,"0 · CRÍTICO        60 · ATENÇÃO        85 · SAUDÁVEL        100",10,F_MONO,
             RGBColor(0xB9,0xAB,0xCF) if dark else MUTED)

def _band(score):
    return ("Saudável",GREEN) if score>=85 else (("Atenção",AMBER) if score>=60 else ("Crítico",RED))

def kpi_card(slide,x,y,w,h,label,value,caption,status=None):
    card=_rect(slide,x,y,w,h,fill=WHITE,line=HAIR,line_w=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    col={"ok":GREEN,"warn":AMBER,"bad":RED,"pri":BRAND}.get(status,TEXT)
    if status: _rect(slide,x+w-0.28,y+0.16,0.13,0.13,fill=col,shape=MSO_SHAPE.OVAL)
    tf=card.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.16); tf.margin_top=Inches(0.14)
    add_text(tf,label.upper(),10,F_MONO,MUTED)
    add_text(tf,value,30,F_DISP,col if status else TEXT,bold=True)
    add_text(tf,caption,11,F_BODY,MUTED)
    return card

def resumo(prs,f):
    s=_blank(prs); section_header(s,"Visão do mês","Resumo executivo")
    score=effective_health(f)
    b=_g(f,"backup",default={}); st=_g(f,"storage",default={}); sec=_g(f,"security",default={}); lc=_g(f,"lifecycle",default={})
    cards=[]
    cards.append(("Postura","%d"%score,"de 100","ok" if score>=85 else ("warn" if score>=60 else "bad")))
    if b.get("runs"): cards.append(("Backup","%s%%"%_pt(b.get("successPct")),"%s execuções"%_pt(b.get("runs")),"ok" if (b.get("successPct") or 0)>=95 else "warn"))
    if st.get("occupancyPct"): cards.append(("Armazenamento","%s%%"%_pt(st.get("occupancyPct")),"%s de %s GB"%(_pt(st.get("cloudGB")),_pt(st.get("contractedGB"))),"warn" if (st.get("occupancyPct") or 0)>=85 else "ok"))
    if sec.get("total"): cards.append(("Alertas","%d"%sec["total"],"%d crítico(s)"%_g(sec,"bySeverity","critical",default=0),"bad" if _g(sec,"bySeverity","critical",default=0) else "warn"))
    if lc.get("total"): cards.append(("Agentes online","%d/%d"%(lc.get("online",0),lc["total"]),"%d offline"%(lc["total"]-lc.get("online",0)),"warn" if lc.get("online",0)<lc["total"] else "ok"))
    cards=cards[:5]
    n=len(cards); gap=0.18; W=11.93; cw=(W-gap*(n-1))/n
    for i,(lab,val,cap,stt) in enumerate(cards):
        kpi_card(s,0.7+i*(cw+gap),1.85,cw,1.5,lab,val,cap,stt)
    # priorities
    pr=_g(f,"analysis","topPriorities",default=[])[:3]
    if pr:
        for i,p in enumerate(pr):
            box=_rect(s,0.7+i*4.04,3.6,3.85,1.0,fill=WHITE,line=HAIR,line_w=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.14)
            add_text(tf,"PRIORIDADE %d"%(i+1),10,F_MONO,BRAND)
            add_text(tf,p,13,F_BODY,TEXT)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def effective_health(f):
    score=int(_g(f,"analysis","healthScore",default=0))
    has_bkp=bool(_g(f,"serviceMap","backup") and _g(f,"backup","runs",default=0))
    if _g(f,"serviceMap","backup") and not has_bkp:
        sp=float(_g(f,"backup","successPct",default=0) or 0)
        score+=int(round(0.5*(100-sp)))
    return max(0,min(100,score))

def cover(prs,f):
    s=_blank(prs)
    _rect(s,0,0,13.333,7.5,fill=INK)  # canvas
    client=_g(f,"meta","client",default="Cliente"); mpt=month_pt(_g(f,"meta","month",default=""))
    sm=_g(f,"serviceMap",default={})
    svcs=[n for n,k in [("Backup","backup"),("EDR","edr"),("VA","va"),("M365","m365")] if sm.get(k)]
    # logo chip (white) with logo image or wordmark
    chip=_rect(s,0.7,0.6,3.2,1.0,fill=WHITE,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    logo=os.path.join(_CONFIG,"contego_logo.jpg")
    if os.path.exists(logo):
        s.shapes.add_picture(logo,Inches(0.95),Inches(0.78),height=Inches(0.64))
    else:
        add_text(chip.text_frame,"Contego Security",18,F_DISP,BRAND,bold=True)
    # title block
    tb=s.shapes.add_textbox(Inches(0.7),Inches(2.4),Inches(9.5),Inches(2.6)).text_frame; tb.word_wrap=True
    add_text(tb,"RELATÓRIO MENSAL DE PROTEÇÃO GERENCIADA",12,F_MONO,RGBColor(0xC9,0xBD,0xDA))
    add_text(tb,client,54,F_DISP,WHITE,bold=True)
    add_text(tb,"Backup, EDR e Microsoft 365 — desempenho e postura do período.",16,F_BODY,RGBColor(0xD9,0xCF,0xE6))
    # meta row
    mt=s.shapes.add_textbox(Inches(0.7),Inches(4.7),Inches(9),Inches(0.7)).text_frame
    add_text(mt,"PERÍODO: %s        SERVIÇOS: %s"%(mpt," · ".join(svcs) or "—"),12,F_MONO,RGBColor(0xB9,0xAB,0xCF))
    # posture card
    score=effective_health(f); band,_c=_band(score)
    _rect(s,0.7,5.55,11.93,1.4,fill=INK2,line=RGBColor(0x5A,0x3A,0x7A),line_w=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pc=s.shapes.add_textbox(Inches(1.0),Inches(5.75),Inches(5),Inches(0.8)).text_frame
    add_text(pc,"POSTURA DE SEGURANÇA",11,F_MONO,RGBColor(0xC9,0xBD,0xDA))
    add_text(pc,"%d / 100  ·  %s"%(score,band),26,F_DISP,WHITE,bold=True)
    posture_meter(s,score,1.0,6.5,11.3,dark=True)
    return s

def table(slide,x,y,w,headers,rows,col_widths=None,status_idx=None):
    nr=len(rows)+1; nc=len(headers)
    gt=slide.shapes.add_table(nr,nc,Inches(x),Inches(y),Inches(w),Inches(0.4+0.42*len(rows))).table
    gt.first_row=False; gt.horz_banding=False
    if col_widths:
        for i,cw in enumerate(col_widths): gt.columns[i].width=Inches(cw)
    for j,htxt in enumerate(headers):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE
        c.text=""; add_text(c.text_frame,htxt.upper(),10,F_MONO,MUTED)
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE; c.text=""
            color=TEXT
            if status_idx is not None and j==status_idx:
                sv=str(val)
                if sv.startswith("✔") or sv=="OK" or sv.lower().startswith("cumpr"): color=GREEN
                elif sv.startswith("Próximo"): color=AMBER
                else: color=RED
            add_text(c.text_frame,str(val),13,F_BODY,color)
    return gt

def load_sla(config_dir):
    p=os.path.join(config_dir,"sla.json")
    try:
        with open(p,encoding="utf-8-sig") as fh: return json.load(fh)
    except Exception: return {"backupSuccessPct":95,"agentsOnlinePct":98,"storageOccupancyMaxPct":85}

def sla_rows(f,sla):
    rows=[]; b=_g(f,"backup",default={}); lc=_g(f,"lifecycle",default={}); st=_g(f,"storage",default={})
    if _g(f,"serviceMap","backup") and b.get("runs"):
        ok=(b.get("successPct") or 0)>=sla.get("backupSuccessPct",95)
        rows.append(["Sucesso de backup","%d%%"%sla.get("backupSuccessPct",95),"%s%%"%_pt(b.get("successPct")),"✔ Cumprido" if ok else "✖ Não cumprido"])
    if lc.get("total"):
        pct=int(round(100.0*lc.get("online",0)/lc["total"])); ok=pct>=sla.get("agentsOnlinePct",98)
        rows.append(["Agentes online","%d%%"%sla.get("agentsOnlinePct",98),"%d%%"%pct,"✔ Cumprido" if ok else "✖ Não cumprido"])
    if st.get("occupancyPct"):
        ok=(st.get("occupancyPct") or 0)<=sla.get("storageOccupancyMaxPct",85)
        rows.append(["Ocupação de armazenamento","<= %d%%"%sla.get("storageOccupancyMaxPct",85),"%s%%"%_pt(st.get("occupancyPct")),"✔ Cumprido" if ok else "✖ Não cumprido"])
    return rows

def sla(prs,f,rows):
    s=_blank(prs); section_header(s,"Acordo de nível de serviço","SLA e conformidade")
    table(s,0.7,1.9,11.9,["Métrica","Meta","Atual","Situação"],rows,col_widths=[5.0,2.3,2.3,2.3],status_idx=3)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def donut(slide,x,y,size,segments):
    if sum(s[1] for s in segments)==0:
        tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(size),Inches(0.4))
        add_text(tb.text_frame,"Sem dados no período.",12,F_MONO,MUTED)
        return tb
    cd=CategoryChartData(); cd.categories=[s[0] for s in segments]
    cd.add_series("s",tuple(s[1] for s in segments))
    gf=slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,Inches(x),Inches(y),Inches(size),Inches(size),cd)
    ch=gf.chart; ch.has_title=False; ch.has_legend=True
    ch.legend.position=XL_LEGEND_POSITION.RIGHT; ch.legend.include_in_layout=False
    ch.legend.font.size=Pt(11); ch.legend.font.name=F_BODY
    plot=ch.plots[0]; plot.has_data_labels=False
    ser=plot.series[0]
    for i,seg in enumerate(segments):
        pt=ser.points[i]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb=seg[2]
    return gf

def bars(slide,x,y,w,h,cats):
    if sum(c[1] for c in cats)==0:
        tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(0.4))
        add_text(tb.text_frame,"Sem dados no período.",12,F_MONO,MUTED)
        return tb
    cd=CategoryChartData(); cd.categories=[c[0] for c in cats]; cd.add_series("s",tuple(c[1] for c in cats))
    gf=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(x),Inches(y),Inches(w),Inches(h),cd)
    ch=gf.chart; ch.has_title=False; ch.has_legend=False
    plot=ch.plots[0]; plot.vary_by_categories=True; plot.gap_width=80; plot.has_data_labels=True
    plot.data_labels.font.size=Pt(12); plot.data_labels.font.name=F_DISP
    ser=plot.series[0]
    for i,c in enumerate(cats):
        pt=ser.points[i]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb=c[2]
    return gf

def seguranca(prs,f):
    s=_blank(prs); section_header(s,"Detecção e resposta","Postura de segurança & higiene")
    sec=_g(f,"security",default={}); sev=_g(sec,"bySeverity",default={})
    add_text(s.shapes.add_textbox(Inches(0.7),Inches(1.85),Inches(5),Inches(0.3)).text_frame,
             "ALERTAS POR SEVERIDADE",11,F_MONO,MUTED)
    bars(s,0.7,2.2,5.4,3.4,[("Crítico",sev.get("critical",0),RED),("Aviso",sev.get("warning",0),AMBER),("Erro",sev.get("error",0),LIGHT)])
    add_text(s.shapes.add_textbox(Inches(6.6),Inches(1.85),Inches(6),Inches(0.3)).text_frame,
             "PRINCIPAIS TIPOS DE ALERTA",11,F_MONO,MUTED)
    rows=[[t.get("labelPt") or t.get("type"), str(t.get("count",0))] for t in (sec.get("byType") or [])[:6]]
    if rows: table(s,6.6,2.2,6.0,["Tipo","Qtd"],rows,col_widths=[4.8,1.2])
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def va(prs,f):
    s=_blank(prs); section_header(s,"Gestão de vulnerabilidades","Avaliação de vulnerabilidades")
    v=_g(f,"vulnerability",default={})
    kpi_card(s,0.7,2.0,3.8,1.5,"Varreduras","%s"%_pt(v.get("scansRun",0)),"no período","ok")
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def visible_usage(f):
    sm=_g(f,"serviceMap",default={})
    allowed={"device","other"}
    for k in ("backup","edr","va","m365"):
        if sm.get(k): allowed.add(k)
    # consolidate by labelPt
    seen={}
    for r in (_g(f,"serviceUsage",default=[]) or []):
        cat=r.get("category")
        if cat and cat not in allowed: continue
        lab=r.get("labelPt","")
        d=seen.setdefault(lab,{"used":0,"quota":None})
        d["used"]+=float(r.get("used") or 0)
        if r.get("quota"): d["quota"]=(d["quota"] or 0)+float(r.get("quota"))
    out=[]
    for lab,d in seen.items():
        pct=int(round(100*d["used"]/d["quota"])) if d["quota"] else None
        out.append({"labelPt":lab,"used":d["used"],"quota":d["quota"],"pct":pct})
    return out

def m365(prs,f):
    s=_blank(prs); section_header(s,"Nuvem produtiva","Proteção Microsoft 365")
    m=_g(f,"manual",default={})
    _tb=m.get("m365BackupStorageTB")
    _tb_val=("%s TB"%_pt(_tb)) if _tb is not None else "—"
    cards=[("Assentos protegidos","%d/%d"%(m.get("m365SeatsUsed",0),m.get("m365SeatsContracted",0)),"caixas de correio","ok"),
           ("Compartilhados","%d"%m.get("m365SeatsShared",0),"assentos shared",None),
           ("SharePoint","%d"%m.get("m365SharepointSites",0),"sites",None),
           ("Backup M365",_tb_val,"protegido","pri")]
    n=len(cards); gap=0.18; W=11.93; cw=(W-gap*(n-1))/n
    for i,(lab,val,cap,stt) in enumerate(cards):
        kpi_card(s,0.7+i*(cw+gap),2.1,cw,1.6,lab,val,cap,stt)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def servicos(prs,f,rows):
    s=_blank(prs); section_header(s,"Cobertura contratada","Serviços contratados & uso")
    trows=[]
    for r in rows:
        pct="%d%%"%r["pct"] if r["pct"] is not None else "—"
        quota=_pt(r["quota"]) if r["quota"] else "—"
        sit="Próximo do limite" if (r["pct"] or 0)>=90 else "OK"
        trows.append([r["labelPt"],quota,_pt(r["used"]),pct,sit])
    table(s,0.7,1.95,11.9,["Serviço","Cota","Em uso","%","Situação"],trows,col_widths=[5.2,1.7,1.7,1.3,2.0],status_idx=4)
    add_text(s.shapes.add_textbox(Inches(0.7),Inches(6.3),Inches(11),Inches(0.4)).text_frame,
             "Somente serviços contratados são exibidos.",11,F_MONO,MUTED)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def backup(prs,f):
    s=_blank(prs); section_header(s,"Continuidade de dados","Atividade de backup")
    b=_g(f,"backup",default={})
    ok=b.get("ok",0); warn=b.get("warning",0); err=b.get("error",0)
    donut(s,0.7,1.9,3.0,[("Concluídos",ok,GREEN),("Aviso",warn,AMBER),("Erro",err,RED)])
    cen=s.shapes.add_textbox(Inches(0.95),Inches(3.0),Inches(2.5),Inches(0.7)).text_frame
    cen.paragraphs[0].alignment=PP_ALIGN.CENTER
    add_text(cen,"%s%%"%_pt(b.get("successPct")),26,F_DISP,GREEN,bold=True,align=PP_ALIGN.CENTER)
    table(s,0.7,5.0,5.6,["Indicador","Valor"],
          [["Execuções no mês","%s (%s agendadas)"%(_pt(b.get("runs")),_pt(b.get("scheduled")))],
           ["Taxa de sucesso","%s%%"%_pt(b.get("successPct"))],
           ["Volume transferido","~%s GB"%_pt(b.get("transferredGB"))]],col_widths=[3.0,2.6])
    # critical points grouped by cause
    add_text(s.shapes.add_textbox(Inches(6.7),Inches(1.9),Inches(6),Inches(0.3)).text_frame,
             "PONTOS CRÍTICOS — POR CAUSA",11,F_MONO,MUTED)
    y=2.35
    for er in (b.get("errorReasons") or [])[:4]:
        box=_rect(s,6.7,y,5.9,0.8,fill=WHITE,line=RED,line_w=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.16)
        add_text(tf,er.get("plainPt",""),15,F_DISP,TEXT,bold=True)
        add_text(tf,"%d execução(ões) com falha"%er.get("count",0),12,F_BODY,MUTED)
        y+=0.95
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def fora_contrato(prs,f):
    s=_blank(prs); section_header(s,"Atividade fora do escopo","Pontos de atenção — fora do contrato")
    y=1.95
    for o in (_g(f,"outOfContract",default=[]) or []):
        box=_rect(s,0.7,y,11.9,1.25,fill=WHITE,line=AMBER,line_w=2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.12)
        add_text(tf,o.get("labelPt",""),16,F_DISP,TEXT,bold=True)
        add_text(tf,o.get("message",""),13,F_BODY,RGBColor(0x7a,0x5a,0x00))
        m=o.get("metrics") or {}
        if o.get("service")=="backup":
            errs="; ".join("%s (%d)"%(e.get("labelPt"),e.get("count",0)) for e in (m.get("topErrors") or [])[:3])
            add_text(tf,"%s execuções · %s falhas · %s%% sucesso. %s"%(_pt(m.get("runs")),_pt(m.get("error")),_pt(m.get("successPct")),("Erros: "+errs) if errs else ""),12,F_BODY,MUTED)
        y+=1.4
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def plano(prs,f):
    s=_blank(prs); section_header(s,"Próximos passos","Plano de ação")
    pr=_g(f,"analysis","topPriorities",default=[])[:6]
    y=1.95
    for i,p in enumerate(pr,1):
        box=_rect(s,0.7,y,11.9,0.78,fill=WHITE,line=HAIR,line_w=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.2)
        add_text(tf,"%d.  %s"%(i,p),15,F_BODY,TEXT)
        y+=0.9
    if not pr:
        add_text(s.shapes.add_textbox(Inches(0.7),Inches(2),Inches(11),Inches(0.5)).text_frame,
                 "Sem ações prioritárias no período.",14,F_BODY,MUTED)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

def metodologia(prs,f):
    s=_blank(prs); section_header(s,"Como este relatório é feito","Metodologia")
    tf=s.shapes.add_textbox(Inches(0.7),Inches(2.0),Inches(11.9),Inches(4)).text_frame; tf.word_wrap=True
    for line in [
        "Fonte: Acronis Cyber Protect Cloud (acesso somente leitura).",
        "Coleta paginada de atividades, alertas e agentes; unidades de armazenamento em base binária (GB = 1024³ bytes), como no console.",
        "Período: %s.  Coletado em: %s."%(month_pt(_g(f,"meta","month",default="")),_g(f,"meta","generatedAt",default="—")),
        "Seções exibidas conforme os serviços contratados do cliente.",
    ]:
        add_text(tf,line,14,F_BODY,TEXT)
    footer(s,_g(f,"meta","client",default=""),month_pt(_g(f,"meta","month",default="")))
    return s

# build_deck assembled incrementally by later tasks.
_CONFIG="."
def build_deck(f, config_dir):
    global _CONFIG; _CONFIG=config_dir
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    cover(prs,f)
    sm=_g(f,"serviceMap",default={})
    _sla=sla_rows(f,load_sla(config_dir))
    builders=[]
    builders.append(lambda: resumo(prs,f))
    if _sla: builders.append(lambda: sla(prs,f,_sla))
    if sm.get("backup") and _g(f,"backup","runs",default=0): builders.append(lambda: backup(prs,f))
    if sm.get("edr"): builders.append(lambda: seguranca(prs,f))
    if sm.get("va"):  builders.append(lambda: va(prs,f))
    if sm.get("m365"): builders.append(lambda: m365(prs,f))
    _vu=visible_usage(f)
    if _vu: builders.append(lambda: servicos(prs,f,_vu))
    if _g(f,"outOfContract",default=[]): builders.append(lambda: fora_contrato(prs,f))
    builders.append(lambda: plano(prs,f))
    builders.append(lambda: metodologia(prs,f))
    _SEC["n"]=0; _SEC["total"]=len(builders)
    for b in builders: b()
    return prs

def main():
    if len(sys.argv)<3:
        print("usage: build_pptx.py <facts.json> <out_dir> [config_dir]"); return 2
    facts_path=sys.argv[1]; out_dir=sys.argv[2]
    config_dir=sys.argv[3] if len(sys.argv)>3 else os.path.join(os.path.dirname(os.path.abspath(__file__)),"config")
    f=load_facts(facts_path)
    prs=build_deck(f,config_dir)
    slug=slug_from_path(facts_path); month=_g(f,"meta","month",default="")
    out=os.path.join(out_dir,"RELATORIO_MSSP_%s_%s.pptx"%(slug,month))
    prs.save(out); print("pptx -> %s"%out); return 0

if __name__=="__main__":
    sys.exit(main())
